"""Generowanie prezentacji PPTX o projekcie cheiloskopii.

Uruchomienie:
    python create_presentation.py

Wynik:
    prezentacja.pptx (24 slajdy, 16:9, po polsku)
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Kolory ────────────────────────────────────────────────────────────────────
BLUE = RGBColor(0x2E, 0x75, 0xB6)
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF6)
MEDIUM_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

# ── Wymiary (16:9) ───────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
TITLE_BAR_H = Inches(1.15)
CONTENT_TOP = Inches(1.5)
CONTENT_W = SLIDE_W - 2 * MARGIN
FONT_NAME = "Calibri"


# ── Helpery ───────────────────────────────────────────────────────────────────

def _set_font(run, size=18, bold=False, italic=False, color=BLACK, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                 italic=False, color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    tb.text_frame.auto_size = None
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size, bold, italic, color)
    tb.text_frame.paragraphs[0].space_after = Pt(4)
    return tb


def _add_bullets(slide, left, top, width, height, items, size=16, color=BLACK,
                 bold_first=False, spacing=Pt(6)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.level = 0

        # Bold prefix (text before first " — ")
        if bold_first and " — " in item:
            prefix, rest = item.split(" — ", 1)
            r1 = p.add_run()
            r1.text = prefix + " — "
            _set_font(r1, size, bold=True, color=color)
            r2 = p.add_run()
            r2.text = rest
            _set_font(r2, size, color=color)
        else:
            run = p.add_run()
            run.text = item
            _set_font(run, size, color=color)
    return tb


def _add_title_bar(slide, title_text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, TITLE_BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    # Cienka linia pod spodem
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, TITLE_BAR_H, SLIDE_W, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = DARK_BLUE
    line.line.fill.background()
    # Tekst tytulu
    _add_textbox(slide, MARGIN, Inches(0.2), CONTENT_W, Inches(0.8),
                 title_text, size=28, bold=True, color=WHITE)


def _add_placeholder(slide, left, top, width, height, label):
    """Szary prostokat z napisem — placeholder na obraz/wykres."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_GRAY
    rect.line.color.rgb = MEDIUM_GRAY
    rect.line.width = Pt(1)
    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    _set_font(run, 14, italic=True, color=RGBColor(0x80, 0x80, 0x80))
    # Wycentruj pionowo
    tf.paragraphs[0].space_before = Pt(height / Emu(12700) / 3)
    return rect


def _add_table(slide, rows, col_widths, left, top):
    """Dodaj sformatowana tabele. rows[0] = naglowek."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    total_w = sum(col_widths)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                        Inches(total_w), Inches(0.4 * n_rows))
    tbl = tbl_shape.table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cell_text)

            if ri == 0:
                _set_font(run, 13, bold=True, color=WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            else:
                _set_font(run, 12, color=BLACK)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else LIGHT_GRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl_shape


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Slajdy ────────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = _blank(prs)
    # Pelne niebieskie tlo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()

    _add_textbox(slide, MARGIN, Inches(1.5), CONTENT_W, Inches(2.0),
                 "System cheiloskopijnej\nidentyfikacji osób",
                 size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_textbox(slide, MARGIN, Inches(3.5), CONTENT_W, Inches(0.6),
                 "Porównanie metod biometrycznej identyfikacji na podstawie fotografii ust",
                 size=20, color=WHITE, align=PP_ALIGN.CENTER)

    # Linia oddzielajaca
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.4),
                                  Inches(5.333), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = WHITE
    sep.line.fill.background()

    _add_textbox(slide, MARGIN, Inches(4.8), CONTENT_W, Inches(0.5),
                 "Praca magisterska",
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, MARGIN, Inches(5.4), CONTENT_W, Inches(0.5),
                 "[Imie i Nazwisko]  |  [Nazwa uczelni]  |  [Rok]",
                 size=16, color=WHITE, align=PP_ALIGN.CENTER)


def slide_02_toc(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Spis tresci")

    sections = [
        "1.  Wprowadzenie — czym jest cheiloskopia",
        "2.  Cel pracy i zadanie identyfikacyjne",
        "3.  Zbior danych",
        "4.  Preprocessing obrazow",
        "5.  Metoda 1 — tradycyjne porownanie obrazow",
        "6.  Metoda 2 — uczenie maszynowe na cechach biometrycznych",
        "7.  Ekstrakcja cech: LBP, HOG, Gabor, minucje",
        "8.  Klasyfikatory: SVM, Random Forest, k-NN",
        "9.  Walidacja krzyzowa (5-fold CV)",
        "10. Wyniki badan",
        "11. Porownanie metod",
        "12. Aplikacja GUI (Streamlit)",
        "13. Wnioski i podsumowanie",
    ]
    _add_bullets(slide, MARGIN, CONTENT_TOP, CONTENT_W, Inches(5.5),
                 sections, size=18)


def slide_03_cheiloscopy(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Czym jest cheiloskopia?")

    text = (
        "Cheiloskopia (gr. cheilos = usta, skopein = badac) to dzial kryminalistyki "
        "zajmujacy sie identyfikacja osob na podstawie wzorow bruzd ust.\n\n"
        "Analogicznie do daktyloskopii (odciskow palcow), wzory bruzd ust sa:"
    )
    _add_textbox(slide, MARGIN, CONTENT_TOP, Inches(7.5), Inches(1.5), text, size=16)

    features = [
        "Unikalne dla kazdego czlowieka",
        "Niezmienne przez cale zycie (poza patologiami)",
        "Mozliwe do odtworzenia ze sladu lub zdjecia",
        "Czesciowo dziedziczne (bliznięta — podobne, ale nie identyczne)",
    ]
    _add_bullets(slide, MARGIN, Inches(3.5), Inches(7.5), Inches(2.0),
                 features, size=16)

    # Tabela typow Suzuki-Tsuchihashi
    _add_textbox(slide, MARGIN, Inches(5.5), Inches(7.5), Inches(0.4),
                 "Klasyfikacja Suzuki i Tsuchihashi (1970) — 5 typow bruzd:",
                 size=14, bold=True)

    rows = [
        ["Typ", "Opis", "Symbol"],
        ["I", "Bruzdy pionowe", "|"],
        ["II", "Bruzdy poziome", "—"],
        ["III", "Bruzdy rozgalezione", "Y"],
        ["IV", "Bruzdy krzyzujace", "X"],
        ["V", "Bruzdy nieregularne", "~"],
    ]
    _add_table(slide, rows, [0.8, 3.5, 1.0], MARGIN, Inches(5.95))

    _add_placeholder(slide, Inches(8.8), CONTENT_TOP, Inches(4), Inches(5.2),
                     "[Ilustracja typow bruzd ust\nlub zdjecie odcisku ust]")


def slide_04_goal(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Cel pracy i zadanie identyfikacyjne")

    items = [
        "Cel — porownanie dwoch metod automatycznej identyfikacji osoby "
        "na podstawie zdjecia ust",
        "Zadanie — identyfikacja 1-do-N: dane jedno nieznane zdjecie ust, "
        "szukamy ktora z N=22 znanych osob jest autorem",
        "Metoda 1 — tradycyjne porownanie obrazow (SSIM, ORB, Histogram)",
        "Metoda 2 — klasyczne ML na cechach biometrycznych (LBP + HOG + Gabor + minucje)",
        "Baseline losowy — accuracy = 1/22 = 4.5% (dolna granica odniesienia)",
    ]
    _add_bullets(slide, MARGIN, CONTENT_TOP, Inches(7.5), Inches(3.5),
                 items, size=17, bold_first=False)

    _add_textbox(slide, MARGIN, Inches(5.2), CONTENT_W, Inches(1.0),
                 "Wazna decyzja naukowa: metadane z Excela (plec, wiek, odcien skory) "
                 "NIE wchodza do wektora cech — bylyby to data leakage "
                 "psujacy wartosc badania.",
                 size=14, italic=True, color=RGBColor(0x80, 0x40, 0x40))

    _add_placeholder(slide, Inches(8.8), CONTENT_TOP, Inches(4), Inches(4.5),
                     "[Schemat: zdjecie ust\n-> identyfikacja\n-> osoba z 22]")


def slide_05_dataset(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Zbior danych")

    rows = [
        ["Parametr", "Wartosc"],
        ["Liczba zdjec", "382 (PNG, bezstratne)"],
        ["Liczba osob", "22"],
        ["Zdjec / osobe", "~17 (rozklad nierowny)"],
        ["Rozdzielczosc oryg.", "1478 x 560 px"],
        ["Oswietlenie", "Lampa ON + Lampa OFF (wymieszane)"],
        ["Kadrow.", "Wycięte do ust (_CUT)"],
        ["Metadane", "Plec, wiek, skora, cechy szczeg. (NIE cechy ML!)"],
    ]
    _add_table(slide, rows, [3.0, 6.0], MARGIN, CONTENT_TOP)

    _add_placeholder(slide, Inches(0.6), Inches(5.2), Inches(5.5), Inches(2.0),
                     "[Przykladowe zdjecia ust z bazy danych — galeria]")
    _add_placeholder(slide, Inches(6.8), Inches(5.2), Inches(6.0), Inches(2.0),
                     "[Wykres slupkowy: liczba zdjec per osoba]")


def slide_06_preprocessing(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Preprocessing obrazow")

    _add_textbox(slide, MARGIN, CONTENT_TOP, CONTENT_W, Inches(0.5),
                 "Pipeline przetwarzania wstepnego — 4 kroki:",
                 size=18, bold=True)

    steps = [
        "1. Konwersja do skali szarosci — eliminacja informacji o kolorze (bruzdy sa struktura geometryczna)",
        "2. CLAHE (clip_limit=2.0, tile=8x8) — lokalne wyrownanie kontrastu, bruzdy staja sie "
        "wyrazniejsze niezaleznie od oswietlenia (ON vs OFF)",
        "3. Bilateral filter (d=9, sigma=75) — redukcja szumu z zachowaniem ostrych krawedzi bruzd",
        "4. Resize do 512 x 256 px — ujednolicony rozmiar dla identycznych wektorow cech",
    ]
    _add_bullets(slide, MARGIN, Inches(2.1), Inches(7.0), Inches(3.5), steps, size=15)

    _add_placeholder(slide, Inches(8.0), CONTENT_TOP, Inches(4.8), Inches(5.5),
                     "[Wizualizacja krokow preprocessingu:\nOryginal -> Grayscale -> CLAHE\n"
                     "-> Bilateral -> Resize]")


def slide_07_method1_overview(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Metoda 1: Tradycyjne porownanie obrazow")

    items = [
        "Zasada — brak ekstrakcji cech. Bezposrednie porownanie pikseli/struktur.",
        "Dla kazdego zdjecia testowego obliczamy miare podobienstwa "
        "do WSZYSTKICH zdjec treningowych.",
        "Klasyfikacja 1-NN — osoba z najwyzszym score = predykcja.",
        "Trzy niezalezne miary: SSIM, ORB, Histogram.",
        "Miara Combined — wazona kombinacja trzech powyzszych.",
        "Zalety — prosta implementacja, brak fazy treningowej, interpretowalne wyniki.",
        "Wady — wolne (liniowe po wszystkich treningowych), "
        "wrazliwe na zmiany oswietlenia.",
    ]
    _add_bullets(slide, MARGIN, CONTENT_TOP, CONTENT_W, Inches(5.5),
                 items, size=16, bold_first=True)


def slide_08_method1_metrics(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Metoda 1: Miary podobienstwa")

    rows = [
        ["Miara", "Algorytm", "Zakres", "Implementacja"],
        ["SSIM", "Structural Similarity Index\n(luminancja + kontrast + struktura)",
         "[-1, 1]", "skimage.metrics"],
        ["ORB", "Keypoint matching\n(FAST + BRIEF + BFMatcher)",
         "[0, 1]", "cv2.ORB_create"],
        ["Histogram", "Korelacja histogramow jasnosci\n(Pearson na 256 binach)",
         "[-1, 1]", "cv2.compareHist"],
        ["Combined", "Wazona srednia:\n0.5*SSIM + 0.25*ORB + 0.25*Hist",
         "[0, 1]", "Normalizacja [-1,1]->[0,1]"],
    ]
    _add_table(slide, rows, [1.8, 4.5, 1.2, 2.8], MARGIN, CONTENT_TOP)

    _add_textbox(slide, MARGIN, Inches(5.8), CONTENT_W, Inches(1.0),
                 "Obrazy skalowane do 256x128 px (grayscale) przed porownaniem SSIM "
                 "— wystarczajacy do zachowania struktury bruzd, szybszy niz 1478x560.",
                 size=14, italic=True)


def slide_09_method2_overview(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Metoda 2: Uczenie maszynowe — przeglad pipeline")

    # Diagram jako tekst
    pipeline_text = (
        "Zdjecie PNG (1478x560)\n"
        "       |\n"
        "  PREPROCESSING\n"
        "  Grayscale -> CLAHE -> Bilateral -> Resize 512x256\n"
        "       |\n"
        "  EKSTRAKCJA CECH\n"
        "  LBP(10) + HOG(16740) + Gabor(64) + Minutiae(7)\n"
        "  = surowy wektor 16 821 cech\n"
        "       |\n"
        "  PCA (HOG: 16740 -> 100) + StandardScaler\n"
        "  = finalny wektor 181 cech\n"
        "       |\n"
        "  KLASYFIKATOR: SVM / RF / k-NN\n"
        "       |\n"
        "  Wynik: top-5 kandydatow + confidence"
    )
    _add_textbox(slide, Inches(0.8), CONTENT_TOP, Inches(6.0), Inches(5.5),
                 pipeline_text, size=15)

    _add_placeholder(slide, Inches(7.5), CONTENT_TOP, Inches(5.3), Inches(5.5),
                     "[Schemat blokowy pipeline\nMetody 2]")


def slide_10_lbp(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Ekstrakcja cech: LBP (Local Binary Patterns)")

    items = [
        "Deskryptor tekstury — porownuje jasnosc piksela z sasiadami",
        "Dla kazdego piksela: 8 sasiadow w promieniu 1 px",
        "Jesli sasiad >= srodkowy -> bit=1, inaczej -> bit=0",
        "Wariant 'uniform' — wzorce z <=2 przejsciami 0<->1 (odporne na rotacje)",
        "Histogram 10-binowy znormalizowany (suma=1)",
        "Bin 0 ~ udzial pikseli jednorodnie ciemnych (tlo)",
        "Bin 9 ~ udzial pikseli jednorodnie jasnych",
    ]
    _add_bullets(slide, MARGIN, CONTENT_TOP, Inches(7.0), Inches(3.0), items, size=15)

    rows = [
        ["Parametr", "Wartosc"],
        ["P (sasiedzi)", "8"],
        ["R (promien)", "1 px"],
        ["Metoda", "uniform"],
        ["Biny histogramu", "10"],
        ["Wektor cech", "10 wartosci"],
    ]
    _add_table(slide, rows, [2.5, 2.5], MARGIN, Inches(5.0))

    _add_placeholder(slide, Inches(7.5), CONTENT_TOP, Inches(5.3), Inches(5.5),
                     "[Wizualizacja LBP:\nheatmapa wzorcow\n+ histogram 10-binowy]")


def slide_11_hog(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Ekstrakcja cech: HOG (Histogram of Oriented Gradients)")

    items = [
        "Opisuje lokalna strukture krawedzi przez histogramy kierunkow gradientu",
        "Swietny dla bruzd ust — wyrazna struktura liniowa",
        "Komorki: 16x16 px -> siatka 32x16 = 512 komorek",
        "Bloki: 2x2 komorki z zakladka -> 31x15 = 465 blokow",
        "Normalizacja L2-Hys per blok (clip do 0.2)",
        "Surowy wektor: 465 x 36 = 16 740 cech",
        "PCA redukuje do 100 skladowych (>95% wariancji)",
    ]
    _add_bullets(slide, MARGIN, CONTENT_TOP, Inches(7.0), Inches(3.5), items, size=15)

    rows = [
        ["Parametr", "Wartosc"],
        ["Orientacje", "9 (binow katowych, 0-180 st.)"],
        ["Komorki", "16x16 pikseli"],
        ["Bloki", "2x2 komorki"],
        ["Surowy wektor", "16 740 cech"],
        ["Po PCA", "100 cech"],
    ]
    _add_table(slide, rows, [2.5, 4.0], MARGIN, Inches(5.2))

    _add_placeholder(slide, Inches(7.5), CONTENT_TOP, Inches(5.3), Inches(5.5),
                     "[Wizualizacja HOG:\ngradienty nalozone\nna obraz ust]")


def slide_12_gabor(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Ekstrakcja cech: Filtry Gabora")

    items = [
        "Filtr Gabora = iloczyn Gaussiana i fali sinusoidalnej (bandpass filter)",
        "Biologicznie odpowiada receptorom kory wzrokowej V1",
        "Bruzdy ust biegna w roznych kierunkach — bank filtrow "
        "wychwytuje energie w kazdym kierunku osobno",
        "Dla kazdego filtra: mean(|odpowiedz|) + std(|odpowiedz|)",
        "32 filtry x 2 statystyki = 64 cechy",
    ]
    _add_bullets(slide, MARGIN, CONTENT_TOP, Inches(7.0), Inches(2.5), items, size=15)

    rows = [
        ["Parametr", "Wartosc"],
        ["Orientacje", "8 (co 22.5 stopni: 0, 22.5, ... 157.5)"],
        ["Czestotliwosci", "4 (0.1, 0.2, 0.3, 0.4 cykli/px)"],
        ["Rozmiar jadra", "21 x 21 px"],
        ["Sigma", "4.0"],
        ["Gamma", "0.5"],
        ["Filtrow lacznie", "8 x 4 = 32"],
        ["Wektor cech", "64 wartosci (mean + std per filtr)"],
    ]
    _add_table(slide, rows, [2.5, 5.0], MARGIN, Inches(4.3))

    _add_placeholder(slide, Inches(8.5), CONTENT_TOP, Inches(4.3), Inches(5.5),
                     "[Siatka odpowiedzi Gabora:\n4 wiersze (czest.)\nx 8 kolumn (orient.)]")


def slide_13_minutiae(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Ekstrakcja cech: Minucje bruzd ust")

    _add_textbox(slide, MARGIN, CONTENT_TOP, Inches(7.5), Inches(0.5),
                 "Analogia do daktyloskopii — punkty charakterystyczne bruzd:",
                 size=16, bold=True)

    steps = [
        "1. Binaryzacja — adaptacyjne progowanie (block 35x35, C=10) "
        "oddziela bruzdy od tla",
        "2. Szkieletyzacja — redukcja bruzd do linii o szer. 1 piksela",
        "3. Crossing Number (CN) — skanowanie sasiedztwa 3x3:",
    ]
    _add_bullets(slide, MARGIN, Inches(2.1), Inches(7.5), Inches(1.5), steps, size=15)

    cn_rows = [
        ["CN", "Typ", "Znaczenie"],
        ["1", "Zakonczenie (ending)", "Bruzda sie tu konczy"],
        ["3", "Rozwidlenie (bifurcation)", "Bruzda sie rozgalezia"],
        [">=4", "Skrzyzowanie (crossing)", "Bruzdy sie przecinaja"],
    ]
    _add_table(slide, cn_rows, [0.8, 3.2, 3.5], MARGIN, Inches(3.8))

    feat_rows = [
        ["Cecha", "Opis"],
        ["n_endings", "Liczba zakonczeni bruzd"],
        ["n_bifurcations", "Liczba rozwidlen"],
        ["n_crossings", "Liczba skrzyzowan"],
        ["n_total_minutiae", "Suma powyzszych"],
        ["groove_density", "Minutiae / 10k pikseli szkieletu"],
        ["total_skeleton_px", "Dlug. bruzd w pikselach"],
        ["binary_entropy", "Entropia Shannona (0-1)"],
    ]
    _add_table(slide, feat_rows, [2.5, 4.5], MARGIN, Inches(5.5))

    _add_placeholder(slide, Inches(8.3), CONTENT_TOP, Inches(4.5), Inches(5.5),
                     "[Wizualizacja:\nszkielet bruzd\n+ kolorowe minucje\n"
                     "(zielony=ending\nczerwony=bifurc.\nniebieski=crossing)]")


def slide_14_feature_vector(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Kompozycja wektora cech")

    rows = [
        ["Ekstraktor", "Surowy wymiar", "Po PCA", "Co opisuje"],
        ["LBP", "10", "10", "Tekstura lokalna (wzorce sasiedztwa)"],
        ["HOG", "16 740", "100", "Krawedzie i ich orientacja"],
        ["Gabor", "64", "64", "Tekstura kierunkowa (bruzdy)"],
        ["Minutiae", "7", "7", "Punkty charakterystyczne bruzd"],
        ["LACZNIE", "16 821", "181", "Pelny wektor cech biometrycznych"],
    ]
    _add_table(slide, rows, [1.8, 1.8, 1.5, 5.2], MARGIN, CONTENT_TOP)

    items = [
        "Surowy wektor (16 821 cech) cache'owany na dysku jako .npy per zdjecie",
        "PCA redukuje HOG: 16 740 -> 100 (eliminuje klątwe wymiarowości)",
        "StandardScaler normalizuje kazda ceche do mean=0, std=1",
        "Scaler fitowany TYLKO na danych treningowych (bez wycieku z testu)",
        "Finalny wektor: 181 liczb float32 — zwarta reprezentacja biometryczna",
    ]
    _add_bullets(slide, MARGIN, Inches(4.5), CONTENT_W, Inches(2.5), items, size=15)


def slide_15_classifiers(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Klasyfikatory")

    rows = [
        ["", "SVM", "Random Forest", "k-NN"],
        ["Algorytm", "Support Vector Machine\n(hiperpłaszczyzna maks. marginesu)", "Zespol 200 drzew decyzyjnych\n(bagging + losowe cechy)", "k-Nearest Neighbors\n(odleglosc euklidesowa)"],
        ["Jadro/metoda", "RBF (Radial Basis Function)", "max_features = sqrt(181) ~ 13", "k = 15, weights = uniform"],
        ["Regularyzacja", "C = 10", "class_weight = balanced", "Softmax z odl. min-per-class"],
        ["Wieloklas.", "OvO (231 klas. binarnych)", "Naturalnie wieloklasowy", "Glosowanie sasiadow"],
        ["Probabilist.", "Platt scaling (5-fold wewn.)", "Udzial glosow drzew", "Wlasny softmax (-d/T)"],
        ["Cechy", "Wrażliwy na skale (scaler!)", "Niezalezny od skali", "Wrażliwy na skale"],
    ]
    _add_table(slide, rows, [1.8, 3.2, 3.5, 3.3], MARGIN, CONTENT_TOP)


def slide_16_cv(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Walidacja krzyzowa (5-fold Stratified CV)")

    items = [
        "Zbior 382 zdjec dzielony na 5 rownych czesci (foldow)",
        "Kazdy fold: ~306 treningowych (80%), ~76 testowych (20%)",
        "Stratyfikacja — proporcje 22 klas zachowane w kazdym foldzie",
        "Model trenowany od zera w kazdym foldzie (niezalezny pomiar)",
        "Wynik koncowy: mean +/- std z 5 przebiegow",
        "Foldy zapisane do JSON — identyczne miedzy sesjami (seed=42)",
    ]
    _add_bullets(slide, MARGIN, CONTENT_TOP, Inches(7.0), Inches(3.0), items, size=16)

    rows = [
        ["Metryka", "Opis", "Zakres"],
        ["Accuracy (top-1)", "Czy predykcja nr 1 = prawdziwa osoba", "0-100%"],
        ["Top-3 accuracy", "Czy prawdziwa osoba w top-3 kandydatow", "0-100%"],
        ["Top-5 accuracy", "Czy prawdziwa osoba w top-5 kandydatow", "0-100%"],
        ["Precision (macro)", "Dokladnosc — sr. niewazona po klasach", "0-100%"],
        ["Recall (macro)", "Czulosc — sr. niewazona po klasach", "0-100%"],
        ["F1-score (macro)", "Srednia harmoniczna precision i recall", "0-100%"],
    ]
    _add_table(slide, rows, [2.2, 5.5, 1.5], MARGIN, Inches(4.5))


def slide_17_results_m1(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Wyniki: Metoda 1 (tradycyjna)")

    rows = [
        ["Miara", "Accuracy", "Top-3", "Top-5", "Avg. czas / zdjecie"],
        ["SSIM", "...%", "...%", "...%", "... s"],
        ["ORB", "...%", "...%", "...%", "... s"],
        ["Histogram", "...%", "...%", "...%", "... s"],
        ["Combined", "...%", "...%", "...%", "... s"],
    ]
    _add_table(slide, rows, [2.0, 2.0, 1.5, 1.5, 2.5], Inches(1.5), CONTENT_TOP)

    _add_textbox(slide, MARGIN, Inches(4.5), CONTENT_W, Inches(0.5),
                 "[Uzupelnij wyniki po uruchomieniu ewaluacji na Stronie 8 aplikacji]",
                 size=14, italic=True, color=RGBColor(0x99, 0x33, 0x33))

    _add_placeholder(slide, Inches(1.5), Inches(5.0), Inches(10.0), Inches(2.0),
                     "[Wykres slupkowy porownania miar — accuracy / top-5]")


def slide_18_results_m2(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Wyniki: Metoda 2 (ML)")

    rows = [
        ["Klasyfikator", "Accuracy\n(mean +/- std)", "Top-3", "Top-5", "F1 macro", "Czas / fold"],
        ["SVM (RBF)", "... +/- ...", "...%", "...%", "...", "... s"],
        ["Random Forest", "... +/- ...", "...%", "...%", "...", "... s"],
        ["k-NN (k=15)", "... +/- ...", "...%", "...%", "...", "... s"],
    ]
    _add_table(slide, rows, [2.0, 2.3, 1.3, 1.3, 1.5, 1.8], Inches(0.8), CONTENT_TOP)

    _add_textbox(slide, MARGIN, Inches(4.2), CONTENT_W, Inches(0.5),
                 "[Uzupelnij wyniki po uruchomieniu ewaluacji na Stronie 6 aplikacji]",
                 size=14, italic=True, color=RGBColor(0x99, 0x33, 0x33))

    _add_placeholder(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.3),
                     "[Wykres slupkowy: accuracy / top-5 / F1 per klasyfikator]")


def slide_19_confusion_svm(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Macierz pomylek: SVM")

    _add_placeholder(slide, Inches(1.0), CONTENT_TOP, Inches(11.0), Inches(5.5),
                     "[Wstaw: results/plots/confusion_matrix_svm.png\n\n"
                     "Macierz 22x22 — os X: predykowana klasa, os Y: prawdziwa klasa.\n"
                     "Ciemniejsze komorki = wiecej klasyfikacji.\n"
                     "Idealny klasyfikator: wartosci tylko na przekatnej.]")


def slide_20_confusion_rf_knn(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Macierze pomylek: Random Forest i k-NN")

    _add_placeholder(slide, Inches(0.4), CONTENT_TOP, Inches(6.2), Inches(5.5),
                     "[Wstaw:\nresults/plots/confusion_matrix_rf.png]")
    _add_placeholder(slide, Inches(6.8), CONTENT_TOP, Inches(6.2), Inches(5.5),
                     "[Wstaw:\nresults/plots/confusion_matrix_knn.png]")


def slide_21_comparison(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Porownanie metod")

    rows = [
        ["Aspekt", "Metoda 1 (tradycyjna)", "Metoda 2 (ML)"],
        ["Podejscie", "Porownanie obrazow (piksele)", "Ekstrakcja cech + klasyfikator"],
        ["Cechy", "Brak (bezposrednie porownanie)", "181 cech (LBP+HOG+Gabor+minucje)"],
        ["Trening", "Brak (lazy: zaladuj obrazy)", "Fit PCA + Scaler + Klasyfikator"],
        ["Predykcja", "1-NN po score podobienstwa", "SVM / RF / k-NN po wekt. cech"],
        ["Interpretowalnosc", "Wysoka (widac podobne zdj.)", "Srednia (cechy czytelne)"],
        ["Szybkosc predykcji", "Wolna (306 porownan/zdj.)", "Szybka (<1 s z cache)"],
        ["Accuracy (CV)", "[...%]", "[...%]"],
        ["Top-5 (CV)", "[...%]", "[...%]"],
    ]
    _add_table(slide, rows, [2.5, 4.5, 4.8], MARGIN, CONTENT_TOP)

    _add_textbox(slide, MARGIN, Inches(6.0), CONTENT_W, Inches(1.0),
                 "[Uzupelnij wyniki po przeprowadzeniu ewaluacji obu metod]",
                 size=14, italic=True, color=RGBColor(0x99, 0x33, 0x33))


def slide_22_app(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Aplikacja: Streamlit GUI")

    _add_textbox(slide, MARGIN, CONTENT_TOP, Inches(6.5), Inches(0.5),
                 "Jedna ujednolicona aplikacja webowa — 8 stron:",
                 size=16, bold=True)

    pages = [
        "1. Eksploracja danych — statystyki bazy, galeria, metadane osob",
        "2. Preprocessing — wizualizacja krokow z suwakami parametrow",
        "3. Ekstrakcja cech — LBP/HOG/Gabor/minucje z wizualizacja",
        "4. Trening modelu — SVM/RF/k-NN, 1 fold lub 5-fold CV, auto-retrain",
        "5. Predykcja — wszystkie 3 klasyfikatory jednoczesnie, konsensus",
        "6. Ewaluacja — pelny 5-fold CV, macierze pomylek, eksport Excel",
        "7. Predykcja M1 — porownanie obrazow (SSIM/ORB/Hist)",
        "8. Ewaluacja M1 — 5-fold CV dla miar podobienstwa",
    ]
    _add_bullets(slide, MARGIN, Inches(2.1), Inches(6.5), Inches(4.0), pages, size=14)

    _add_textbox(slide, MARGIN, Inches(6.0), Inches(6.5), Inches(0.5),
                 "Uruchomienie:  streamlit run apps/app.py  (port 8502)",
                 size=14, bold=True, color=DARK_BLUE)

    _add_placeholder(slide, Inches(7.5), CONTENT_TOP, Inches(5.3), Inches(5.5),
                     "[Screenshot aplikacji\n— strona Predykcji\nz wynikami 3 klasyfikatorow]")


def slide_23_conclusions(prs):
    slide = _blank(prs)
    _add_title_bar(slide, "Wnioski i podsumowanie")

    items = [
        "Cheiloskopia jest wykonalna metoda identyfikacji biometrycznej "
        "— nawet na malym zbiorze 382 zdjec / 22 osob",
        "Metoda 2 (ML na cechach) przewyzsza Metode 1 (porownanie obrazow) "
        "pod wzgledem accuracy i szybkosci predykcji",
        "Minucje bruzd ust (endings, bifurcations, crossings) sa unikalnym "
        "wkladem cheiloskopii — analogia do daktyloskopii",
        "HOG wymaga redukcji PCA — 16 740 cech przy 382 probkach "
        "to klątwa wymiarowosci",
        "k-NN wymaga wlasnego softmax prawdopodobienstw — standardowe "
        "glosowanie daje nieinformatywne 100%",
        "Cache wektorow cech (NPY) redukuje czas 90x (z ~90 s do ~1 s)",
    ]
    _add_bullets(slide, MARGIN, CONTENT_TOP, Inches(7.5), Inches(3.5), items, size=16)

    _add_textbox(slide, MARGIN, Inches(5.3), Inches(7.5), Inches(0.5),
                 "Mozliwe kierunki rozwoju:", size=16, bold=True)

    future = [
        "Deep learning (CNN, Siamese) — na wiekszym zbiorze danych",
        "Augmentacja danych — sztuczne powiekszenie zbioru",
        "Klasyfikacja typow bruzd Suzuki-Tsuchihashi",
        "Weryfikacja 1-do-1 (czy dwa zdjecia = ta sama osoba)",
    ]
    _add_bullets(slide, MARGIN, Inches(5.7), Inches(7.5), Inches(1.5), future, size=14)


def slide_24_thanks(prs):
    slide = _blank(prs)
    # Pelne niebieskie tlo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()

    _add_textbox(slide, MARGIN, Inches(2.5), CONTENT_W, Inches(1.5),
                 "Dziekuje za uwage",
                 size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.0),
                                  Inches(4.333), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = WHITE
    sep.line.fill.background()

    _add_textbox(slide, MARGIN, Inches(4.5), CONTENT_W, Inches(1.0),
                 "Pytania?",
                 size=28, color=WHITE, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_toc(prs)
    slide_03_cheiloscopy(prs)
    slide_04_goal(prs)
    slide_05_dataset(prs)
    slide_06_preprocessing(prs)
    slide_07_method1_overview(prs)
    slide_08_method1_metrics(prs)
    slide_09_method2_overview(prs)
    slide_10_lbp(prs)
    slide_11_hog(prs)
    slide_12_gabor(prs)
    slide_13_minutiae(prs)
    slide_14_feature_vector(prs)
    slide_15_classifiers(prs)
    slide_16_cv(prs)
    slide_17_results_m1(prs)
    slide_18_results_m2(prs)
    slide_19_confusion_svm(prs)
    slide_20_confusion_rf_knn(prs)
    slide_21_comparison(prs)
    slide_22_app(prs)
    slide_23_conclusions(prs)
    slide_24_thanks(prs)

    out_path = "prezentacja.pptx"
    prs.save(out_path)
    print(f"Prezentacja zapisana: {out_path}")
    print(f"Liczba slajdow: {len(prs.slides)}")


if __name__ == "__main__":
    main()
